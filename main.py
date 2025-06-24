import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm
from io import BytesIO
from PIL import Image
import os

# Streamlit UI
st.title("PDF to PPTX Converter(py 팬더)")
st.write("Upload a PDF file to convert each page to a slide in a PPTX file.")

# File uploader
uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

# Convert PDF to PPTX
def convert_pdf_to_pptx(pdf_data, output_filename):
    pdf_document = fitz.open("pdf", pdf_data)
    presentation = Presentation()

    for page_num in tqdm(range(len(pdf_document)), desc="Converting PDF to PPTX"):
        page = pdf_document.load_page(page_num)
        
        # 페이지 크기 가져오기 (인치 단위로 변환)
        page_width = Inches(page.rect.width / 72)  # 1인치 = 72pt
        page_height = Inches(page.rect.height / 72)
        
        # 슬라이드 크기 설정
        presentation.slide_width = page_width
        presentation.slide_height = page_height
        
        # PDF 페이지를 이미지로 렌더링
        pix = page.get_pixmap()
        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # 슬라이드 추가 및 이미지 배치
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # 빈 슬라이드 레이아웃
        
        image_width, image_height = image.size
        aspect_ratio = image_width / image_height

        if page_width / page_height > aspect_ratio:
            new_height = page_height
            new_width = new_height * aspect_ratio
        else:
            new_width = page_width
            new_height = new_width / aspect_ratio

        left = (page_width - new_width) / 2
        top = (page_height - new_height) / 2
        
        # 이미지를 슬라이드에 추가
        image_bytes = BytesIO()
        image.save(image_bytes, format="PNG")
        image_bytes.seek(0)
        slide.shapes.add_picture(image_bytes, left, top, width=new_width, height=new_height)
    
    # 프레젠테이션을 BytesIO에 저장하여 반환
    pptx_data = BytesIO()
    presentation.save(pptx_data)
    pptx_data.seek(0)
    return pptx_data

# Process PDF and provide download link
if uploaded_file is not None:
    # PDF 파일 이름에서 확장자를 제외하고 PPTX 파일 이름 생성
    output_filename = os.path.splitext(uploaded_file.name)[0] + ".pptx"
    
    st.write("Converting PDF to PPTX, please wait...")
    pptx_data = convert_pdf_to_pptx(uploaded_file.read(), output_filename)
    st.success("Conversion completed!")

    # Provide download link for PPTX
    st.download_button(
        label="Download PPTX file",
        data=pptx_data,
        file_name=output_filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
